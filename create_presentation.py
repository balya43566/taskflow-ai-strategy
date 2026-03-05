from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Modern color palette ───────────────────────────────────────────────────────
NEAR_BLACK  = RGBColor(0x0D, 0x1B, 0x2A)
NAVY        = RGBColor(0x1B, 0x3A, 0x6B)
BLUE        = RGBColor(0x2D, 0x6A, 0xCF)
BLUE_PALE   = RGBColor(0xEB, 0xF3, 0xFF)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_PALE    = RGBColor(0xFF, 0xED, 0xED)
AMBER       = RGBColor(0xC0, 0x6B, 0x00)
AMBER_PALE  = RGBColor(0xFF, 0xF3, 0xD6)
GREEN       = RGBColor(0x05, 0x7A, 0x55)
GREEN_PALE  = RGBColor(0xD6, 0xF5, 0xEB)
GRAY_100    = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_200    = RGBColor(0xE2, 0xE8, 0xF0)
GRAY_400    = RGBColor(0x94, 0xA3, 0xB8)
GRAY_600    = RGBColor(0x47, 0x55, 0x6E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Core primitives ────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill, border=None, border_w=0.75):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(border_w)
    else:
        sh.line.fill.background()
    return sh

def oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h,
        size=16, bold=False, italic=False, color=NEAR_BLACK,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box

def bullets(slide, items, l, t, w, h, size=14, color=NEAR_BLACK,
            dot_color=None, spacing=6):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"

def pagenum(slide, n, light=False):
    c = GRAY_400 if not light else RGBColor(0x60, 0x80, 0xB0)
    txt(slide, str(n), 12.7, 7.15, 0.55, 0.3, size=9, color=c, align=PP_ALIGN.RIGHT)

# ── Slide chrome ───────────────────────────────────────────────────────────────

def content_slide(slide, title, section, accent=BLUE):
    """White background + left accent strip + title block."""
    rect(slide, 0, 0, 13.33, 7.5, WHITE)           # white bg
    rect(slide, 0, 0, 0.08, 7.5, accent)            # left strip
    # bottom footer line
    rect(slide, 0, 7.38, 13.33, 0.12, GRAY_100)
    txt(slide, "TaskFlow  ·  Confidential  ·  March 2025",
        0.25, 7.39, 8, 0.12, size=7, color=GRAY_400)
    # section label
    txt(slide, section.upper(), 0.28, 0.2, 10, 0.28,
        size=8, bold=True, color=accent)
    # title
    txt(slide, title, 0.28, 0.44, 12.9, 0.75,
        size=30, bold=True, color=NEAR_BLACK)
    # divider
    rect(slide, 0.28, 1.25, 12.8, 0.025, GRAY_200)

def card(slide, l, t, w, h, fill=WHITE, accent_color=None, accent_w=0.06):
    """Card with optional left color accent and subtle border."""
    rect(slide, l, t, w, h, fill, border=GRAY_200, border_w=0.5)
    if accent_color:
        rect(slide, l, t, accent_w, h, accent_color)

def dark_slide(slide, bg=NAVY):
    rect(slide, 0, 0, 13.33, 7.5, bg)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s, NEAR_BLACK)

# Right panel – subtle lighter block
rect(s, 8.8, 0, 4.53, 7.5, NAVY)
rect(s, 8.75, 0, 0.04, 7.5, BLUE)   # separator

# Confidential badge
rect(s, 0.45, 0.48, 2.1, 0.36, RED)
txt(s, "CONFIDENTIAL  ·  EXEC TEAM", 0.5, 0.5, 2.0, 0.3,
    size=8, bold=True, color=WHITE)

# Main title
txt(s, "TaskFlow", 0.45, 1.05, 8, 0.95,
    size=58, bold=True, color=WHITE)
txt(s, "AI Competitive\nResponse Strategy", 0.45, 1.95, 8, 1.4,
    size=34, color=RGBColor(0x7B, 0xB3, 0xF0))

# Thin accent line
rect(s, 0.45, 3.5, 4.5, 0.04, BLUE)

txt(s, "Market landscape  ·  Threat assessment  ·  Recommended actions",
    0.45, 3.65, 8.1, 0.4, size=13, italic=True, color=GRAY_400)
txt(s, "Prepared for Sarah Chen, CEO  ·  March 2025",
    0.45, 4.15, 8.1, 0.35, size=12, color=GRAY_400)

# Right panel — threat summary
txt(s, "COMPETITOR THREAT LEVELS", 9.1, 0.45, 3.9, 0.35,
    size=8, bold=True, color=GRAY_400)

for i, (label, sub, clr) in enumerate([
    ("VERY HIGH", "ClickUp Brain", RED),
    ("HIGH", "Asana · Notion · Monday", RED),
    ("MEDIUM", "Linear", AMBER),
]):
    y = 0.95 + i * 1.9
    rect(s, 9.1, y, 3.85, 1.7, clr)
    txt(s, label, 9.25, y + 0.18, 3.5, 0.45,
        size=18, bold=True, color=WHITE)
    txt(s, sub, 9.25, y + 0.62, 3.5, 0.9,
        size=12, color=WHITE, italic=True)

pagenum(s, 1, light=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TASKFLOW LOGO / BRAND
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)

# ── Icon: navy square with task-row motif ──────────────────────────────────
IX, IY, IS = 4.42, 1.65, 2.1   # icon x, y, size (square)
rect(s, IX, IY, IS, IS, NEAR_BLACK)

# Three task rows inside the icon
for row in range(3):
    ry = IY + 0.38 + row * 0.57
    # Checkbox square
    rect(s, IX + 0.22, ry, 0.28, 0.28, WHITE)
    # On the top row: filled checkbox with a checkmark feel
    if row == 0:
        rect(s, IX + 0.22, ry, 0.28, 0.28, BLUE)
        # tick: two small rects forming a check
        rect(s, IX + 0.27, ry + 0.14, 0.06, 0.1,  WHITE)
        rect(s, IX + 0.30, ry + 0.08, 0.14, 0.06, WHITE)
    # Task line
    line_w = 1.38 if row < 2 else 0.9   # shorter last line for variety
    rect(s, IX + 0.62, ry + 0.08, line_w, 0.12, WHITE if row < 2 else GRAY_400)

# Blue accent corner triangle (bottom-right of icon)
rect(s, IX + IS - 0.42, IY + IS - 0.42, 0.42, 0.42, BLUE)

# ── Wordmark ───────────────────────────────────────────────────────────────
# "Task" dark, "Flow" blue — side by side
WORD_Y = IY + IS + 0.32
txt(s, "Task", 3.62, WORD_Y, 2.6, 0.85,
    size=52, bold=True, color=NEAR_BLACK, align=PP_ALIGN.RIGHT)
txt(s, "Flow", 6.28, WORD_Y, 3.2, 0.85,
    size=52, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

# ── Tagline ────────────────────────────────────────────────────────────────
rect(s, 5.42, WORD_Y + 0.92, 2.5, 0.03, GRAY_200)
txt(s, "Work that flows.", 0.5, WORD_Y + 1.05, 12.33, 0.4,
    size=14, italic=True, color=GRAY_400, align=PP_ALIGN.CENTER)

# ── Subtle bottom brand strip ──────────────────────────────────────────────
rect(s, 0, 7.3, 13.33, 0.2, NEAR_BLACK)
txt(s, "AI Competitive Response Strategy  ·  March 2025  ·  Confidential",
    0, 7.32, 13.33, 0.16, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — EXECUTIVE SUMMARY  (was slide 2)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Executive Summary", "Overview", BLUE)
pagenum(s, 3)

# Three stat cards
for i, (stat, label, clr, pale) in enumerate([
    ("5 of 5",      "competitors have\nsome form of AI chat", BLUE,  BLUE_PALE),
    ("30–90",       "days until AI-driven\nchurn risk peaks",  RED,   RED_PALE),
    ("$200–280/mo", "what competitors charge\na 10-person team", AMBER, AMBER_PALE),
]):
    x = 0.28 + i * 4.3
    rect(s, x, 1.38, 4.0, 2.1, pale, border=GRAY_200, border_w=0.5)
    rect(s, x, 1.38, 0.06, 2.1, clr)
    txt(s, stat, x + 0.22, 1.5, 3.7, 0.9,
        size=38, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, label, x + 0.1, 2.38, 3.8, 0.85,
        size=13, color=GRAY_600, align=PP_ALIGN.CENTER)

# Bottom section
txt(s, "Key Takeaways", 0.28, 3.7, 4, 0.35, size=11, bold=True, color=BLUE)
rect(s, 0.28, 4.05, 12.8, 0.025, GRAY_200)

items = [
    ("AI is now a purchase criterion", "Buyers ask 'what AI?' before 'what price?' — this shift happened in 2024.",
     BLUE),
    ("The window is not closed", "Churn happens at renewal. TaskFlow has 30–90 days to communicate a credible AI roadmap.",
     AMBER),
    ("ClickUp Brain is threat #1", "Most mature 'chat with your work' product. Direct overlap with TaskFlow's core use case.",
     RED),
    ("A pricing wedge exists", "Every competitor charges $15–28/user extra for AI. TaskFlow can undercut on price.",
     GREEN),
]
for i, (heading, body, clr) in enumerate(items):
    x = 0.28 + (i % 2) * 6.5
    y = 4.2 + (i // 2) * 1.42
    oval(s, x, y + 0.12, 0.22, 0.22, clr)
    txt(s, heading, x + 0.34, y + 0.08, 5.8, 0.3, size=13, bold=True, color=NEAR_BLACK)
    txt(s, body,    x + 0.34, y + 0.4,  5.8, 0.85, size=12, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — COMPETITIVE LANDSCAPE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "The AI Feature Race — Where Every Competitor Stands", "Competitive Landscape", BLUE)
pagenum(s, 4)

tiers = [
    ("TIER 1  —  FULL AI PLATFORM", RED, RED_PALE,
     "ClickUp Brain  ·  Notion AI + Q&A",
     "Cross-workspace knowledge indexing · Conversational AI chat · Autonomous agents · Meeting intelligence"),
    ("TIER 2  —  STRONG / ENTERPRISE", AMBER, AMBER_PALE,
     "Asana AI Teammates  ·  Monday.com AI",
     "Enterprise-grade automations · M365 Copilot integration · Conversational onboarding agents"),
    ("TIER 3  —  NICHE / DEVELOPER", GREEN, GREEN_PALE,
     "Linear + GitHub Copilot",
     "Best-in-class for engineering teams only · No general-purpose AI chat for business users"),
]
for i, (label, clr, pale, tools, desc) in enumerate(tiers):
    y = 1.42 + i * 1.85
    rect(s, 0.28, y, 12.8, 1.7, pale, border=GRAY_200, border_w=0.4)
    rect(s, 0.28, y, 0.08, 1.7, clr)

    # Tier badge
    rect(s, 0.44, y + 0.22, 3.1, 0.36, clr)
    txt(s, label, 0.5, y + 0.24, 3.0, 0.32, size=9, bold=True, color=WHITE)

    txt(s, tools, 3.72, y + 0.2, 9.1, 0.45, size=18, bold=True, color=NEAR_BLACK)
    txt(s, desc,  0.44, y + 0.78, 12.4, 0.75, size=12, italic=True, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FEATURE COMPARISON TABLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "AI Feature Comparison", "Competitive Landscape", BLUE)
pagenum(s, 5)

COL_X = [0.28, 4.18, 5.88, 7.58, 9.28, 10.98]
COL_W = [3.75, 1.55, 1.55, 1.55, 1.55, 1.55]

# Header row
rect(s, 0.28, 1.38, 12.8, 0.52, NEAR_BLACK)
txt(s, "Feature", COL_X[0]+0.15, 1.42, COL_W[0], 0.44,
    size=11, bold=True, color=WHITE)
for col, (cx, cw) in enumerate(zip(COL_X[1:], COL_W[1:])):
    label = ["Asana","Linear","Monday","ClickUp","Notion"][col]
    txt(s, label, cx, 1.42, cw, 0.44, size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

ROWS = [
    ("Chat with tasks/work",
     ("YES","NO","PARTIAL","YES","YES")),
    ("AI automations",
     ("YES","NO","YES","YES","PARTIAL")),
    ("Autonomous AI agents",
     ("YES — β","PARTIAL","YES","YES","YES")),
    ("Meeting intelligence",
     ("NO","NO","NO","YES","YES")),
    ("Cross-platform context",
     ("YES","PARTIAL","PARTIAL","YES","YES")),
    ("Multi-model AI",
     ("NO","NO","NO","YES","NO")),
    ("Enterprise AI security",
     ("YES","NO","YES","YES","YES")),
]

VAL_COLOR = {"YES": GREEN, "PARTIAL": AMBER, "NO": RGBColor(0xCC,0xCC,0xCC),
             "YES — β": AMBER, "YES — β": AMBER}
VAL_PALE  = {"YES": GREEN_PALE, "PARTIAL": AMBER_PALE, "NO": GRAY_100,
             "YES — β": AMBER_PALE}

for r, (feat, vals) in enumerate(ROWS):
    y = 1.9 + r * 0.73
    bg = WHITE if r % 2 == 0 else GRAY_100
    rect(s, 0.28, y, 12.8, 0.73, bg)
    if r % 2 == 0:
        rect(s, 0.28, y, 12.8, 0.73, WHITE)
    else:
        rect(s, 0.28, y, 12.8, 0.73, GRAY_100)

    txt(s, feat, COL_X[0]+0.15, y+0.17, COL_W[0]-0.1, 0.4, size=12, color=NEAR_BLACK)

    for c, val in enumerate(vals):
        cx = COL_X[c+1]
        cw = COL_W[c+1]
        clr = VAL_COLOR.get(val, GRAY_400)
        pale = VAL_PALE.get(val, GRAY_100)
        # Pill background
        rect(s, cx + 0.2, y + 0.18, cw - 0.38, 0.36, pale, border=None)
        display = val.replace("YES — β", "Yes β").replace("YES","Yes").replace("NO","No").replace("PARTIAL","Partial")
        txt(s, display, cx + 0.18, y + 0.19, cw - 0.34, 0.36,
            size=11, bold=(val in ("YES","YES — β")), color=clr,
            align=PP_ALIGN.CENTER)

# Bottom border
rect(s, 0.28, 1.9 + 7*0.73, 12.8, 0.025, GRAY_200)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THREAT ASSESSMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Who Poses the Biggest Risk to TaskFlow?", "Threat Assessment", RED)
pagenum(s, 6)

THREATS = [
    ("ClickUp Brain",      "VERY HIGH", RED,
     "Most comprehensive 'chat with your work' product on the market. Cross-workspace indexing of tasks, docs, comments, and chat. Multi-model AI (GPT-5, Claude, o3). Meeting Notetaker closes the loop from call to task. Direct, head-on overlap with TaskFlow's core use case."),
    ("Notion AI Q&A",      "HIGH",      RED,
     "Best Q&A implementation for knowledge-heavy teams. Notion 3.0 AI Agents execute autonomous multi-step workflows. Cross-platform context from Slack, Google Drive, and Teams. Risk: users consolidate TaskFlow + Notion docs onto a single AI-powered platform."),
    ("Asana AI Teammates", "HIGH",      RED,
     "Enterprise-grade with Microsoft 365 Copilot integration — a massive enterprise advantage. AI Teammates (beta) represent a qualitative new experience. Strong enterprise sales motion already targeting TaskFlow accounts in head-to-head deals."),
    ("Monday.com AI",      "HIGH (enterprise)", AMBER,
     "$1.2B 2025 revenue target with AI central to growth strategy. Monday Expert conversational agent launched March 2025 — compelling onboarding and setup experience. Biggest threat for TaskFlow's enterprise and mid-market customers."),
    ("Linear + Copilot",   "MEDIUM",    AMBER,
     "Engineering-team-only threat via GitHub Copilot agent integration. Powerful for dev orgs already using GitHub. No meaningful risk for general business users or non-technical teams."),
]

for i, (name, level, clr, desc) in enumerate(THREATS):
    y = 1.38 + i * 1.16
    card(s, 0.28, y, 12.8, 1.08, WHITE, accent_color=clr, accent_w=0.08)

    # Threat badge
    rect(s, 0.44, y + 0.28, 2.0, 0.36, clr)
    txt(s, level, 0.46, y + 0.3, 1.96, 0.32, size=9, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, name, 0.44, y + 0.7, 2.1, 0.3, size=12, bold=True, color=NEAR_BLACK)

    txt(s, desc, 2.62, y + 0.14, 10.3, 0.82, size=12, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PERSONAS AT RISK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Which Users Are Most Likely to Churn?", "Threat Assessment", RED)
pagenum(s, 7)

PERSONAS = [
    ("Power Users & Team Admins", "HIGH RISK", RED,
     "Evaluate tools proactively. See competitor demos on LinkedIn and champion switching internally. Acutely feel the pain of manual status reporting and context-switching.",
     "ClickUp Brain · Asana AI Teammates"),
    ("Knowledge Teams\n(PMs, Content, Ops)", "HIGH RISK", RED,
     "'What did we decide about X?' is a daily question. Notion Q&A and ClickUp Brain answer it instantly. These users have the strongest pull toward AI-first alternatives.",
     "Notion Q&A · ClickUp Brain"),
    ("Engineering Teams", "MEDIUM RISK", AMBER,
     "Linear + GitHub Copilot is a purpose-built alternative for dev-centric orgs. CTOs evaluating tools at renewal may prefer Linear's focused, fast experience.",
     "Linear + GitHub Copilot"),
    ("SMB / Solopreneurs", "LOW RISK", GREEN,
     "Price-sensitive and rarely switch for features alone. AI is a 'nice to have' rather than a switching trigger. Least likely to churn in the next 90 days.",
     "Unlikely to switch for AI alone"),
]

for i, (name, risk, clr, desc, competitor) in enumerate(PERSONAS):
    x = 0.28 + (i % 2) * 6.52
    y = 1.38 + (i // 2) * 2.88
    card(s, x, y, 6.18, 2.72, WHITE)
    rect(s, x, y, 6.18, 0.58, clr)
    txt(s, name.replace("\n", "  "), x + 0.18, y + 0.08, 3.8, 0.45,
        size=14, bold=True, color=WHITE)
    # Risk badge
    rect(s, x + 4.45, y + 0.14, 1.55, 0.3, RGBColor(0,0,0,))
    rect(s, x + 4.42, y + 0.12, 1.58, 0.32,
         RGBColor(0xFF,0xFF,0xFF) if clr == GREEN else RGBColor(0x00,0x00,0x00))
    rect(s, x + 4.4, y + 0.1, 1.62, 0.36, WHITE)
    txt(s, risk, x + 4.38, y + 0.12, 1.66, 0.32,
        size=9, bold=True, color=clr, align=PP_ALIGN.CENTER)

    txt(s, desc, x + 0.18, y + 0.72, 5.78, 1.35, size=12, color=GRAY_600)
    rect(s, x + 0.18, y + 2.22, 5.78, 0.02, GRAY_200)
    txt(s, f"Key competitor: {competitor}", x + 0.18, y + 2.3, 5.78, 0.32,
        size=10, italic=True, color=GRAY_400)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TIMELINE PRESSURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "How Much Time Does TaskFlow Have?", "Threat Assessment", RED)
pagenum(s, 8)

# Timeline track
rect(s, 0.5, 3.62, 12.4, 0.14, GRAY_200)
rect(s, 0.5, 3.62, 12.4, 0.14, GRAY_200)

MILESTONES = [
    (0.5,  "NOW",       "Sales needs answers\ntoday — in demo calls",          RED),
    (4.35, "30 DAYS",   "First AI renewal\nconversations begin",                AMBER),
    (8.2,  "90 DAYS",   "First meaningful\nAI-driven churn risk",               RED),
    (11.8, "6 MONTHS",  "AI absence = clear\ncompetitive disadvantage",         RED),
]

for x, label, desc, clr in MILESTONES:
    # Dot on timeline
    oval(s, x - 0.13, 3.52, 0.38, 0.38, clr)
    # Line up
    rect(s, x + 0.06, 2.1, 0.025, 1.45, GRAY_200)
    # Label above
    txt(s, label, x - 0.9, 1.42, 1.95, 0.5,
        size=13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, desc, x - 0.9, 1.92, 1.95, 1.0,
        size=11, color=GRAY_600, align=PP_ALIGN.CENTER)
    txt(s, desc, x - 0.9, 3.97, 1.95, 1.0,
        size=11, color=GRAY_600, align=PP_ALIGN.CENTER)

# Gradient fill of timeline (simulated with rects)
for i, (x_start, x_end, clr) in enumerate([
    (0.5, 4.35, AMBER), (4.35, 8.2, RED), (8.2, 12.5, RED)
]):
    rect(s, x_start, 3.62, x_end - x_start, 0.14, clr)

# Bottom callout
rect(s, 0.5, 5.15, 12.4, 1.65, BLUE_PALE, border=GRAY_200, border_w=0.4)
rect(s, 0.5, 5.15, 0.06, 1.65, BLUE)
txt(s, "The window is not closed — but the clock is running.",
    0.72, 5.3, 12.0, 0.45, size=16, bold=True, color=NEAR_BLACK)
txt(s,
    "Customer churn from AI-related switching typically happens at contract renewal, not immediately. "
    "This means TaskFlow has a meaningful window to communicate a credible AI roadmap — but only if "
    "that communication reaches customers before they hear about competitors from their peers or see a demo.",
    0.72, 5.75, 11.9, 0.95, size=13, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THREE-HORIZON OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Three-Horizon Response Plan", "Response Strategy", BLUE)
pagenum(s, 9)

HORIZONS = [
    ("0–30\nDAYS", "Immediate", RED, [
        "CEO email to customers with AI roadmap preview",
        "Sales battle card with objection handlers",
        "Remove all 'no AI' messaging from website & decks",
    ]),
    ("30–90\nDAYS", "Short-term", AMBER, [
        "Ship AI semantic search (fastest credible feature)",
        "Publish public AI roadmap page",
        "Evaluate Anthropic / OpenAI API partnership",
    ]),
    ("90+\nDAYS", "Long-term", GREEN, [
        "Full conversational AI assistant (RAG-powered)",
        "Natural language AI automations",
        "Meeting intelligence: auto-create tasks from calls",
    ]),
]

for i, (timeframe, label, clr, items) in enumerate(HORIZONS):
    x = 0.28 + i * 4.35
    # Header card
    rect(s, x, 1.38, 4.15, 1.35, clr)
    txt(s, timeframe, x + 0.2, 1.45, 1.4, 1.1,
        size=26, bold=True, color=WHITE)
    txt(s, label.upper(), x + 1.65, 1.55, 2.3, 0.35,
        size=9, bold=True, color=WHITE)

    # Content card
    card(s, x, 2.78, 4.15, 4.42, WHITE)
    for j, item in enumerate(items):
        iy = 3.0 + j * 1.3
        oval(s, x + 0.2, iy + 0.05, 0.28, 0.28, clr)
        txt(s, str(j+1), x + 0.2, iy + 0.04, 0.28, 0.28,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, item, x + 0.6, iy, 3.4, 1.1, size=13, color=NEAR_BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IMMEDIATE ACTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Actions the Team Can Take This Week", "Immediate  ·  0–30 Days", RED)
pagenum(s, 10)

ACTIONS = [
    ("01", "CEO Email to All Customers",
     "Send from Sarah Chen by end of this week.",
     [
         'Subject line: "TaskFlow AI is coming — here\'s what we\'re building"',
         "Acknowledge the AI wave; don't pretend it isn't happening",
         "Share 3 concrete roadmap items with real dates",
         "Offer VIP beta access — turn concern into excitement",
     ]),
    ("02", "Sales Battle Card",
     "One-pager for sales team before next week's pipeline calls.",
     [
         "Objection handler for: ClickUp Brain, Notion Q&A, Asana Teammates",
         "Key counter: competitors charge $200–280/month extra for AI",
         "Bridging answer for 'do you have AI?' — honest + roadmap-forward",
     ]),
    ("03", "Messaging Audit",
     "Immediate website and deck review.",
     [
         "Search for: 'simple', 'focused', 'no bloat', 'just works' — review each",
         "Remove or reframe language that implies intentionally AI-free",
         "What was a differentiator in 2022 is a liability in 2025",
     ]),
]

for i, (num, title, sub, pts) in enumerate(ACTIONS):
    y = 1.38 + i * 1.9
    card(s, 0.28, y, 12.8, 1.78, WHITE, accent_color=RED, accent_w=0.08)
    # Number
    rect(s, 0.44, y + 0.28, 0.65, 0.65, RED)
    txt(s, num, 0.44, y + 0.28, 0.65, 0.65,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    txt(s, title, 1.25, y + 0.14, 6.0, 0.42, size=15, bold=True, color=NEAR_BLACK)
    txt(s, sub,   1.25, y + 0.56, 6.0, 0.3,  size=11, italic=True, color=GRAY_400)
    # Bullets
    for j, pt in enumerate(pts):
        bx = 7.5 + (0 if j == 0 else 0)
        by = y + 0.22 + j * 0.36
        if by < y + 1.68:
            oval(s, 1.28, by + 0.07, 0.13, 0.13, RED)
            txt(s, pt, 1.5, by, 11.4, 0.38, size=11, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SHORT-TERM ACTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "First Shippable AI + Market Signals", "Short-term  ·  30–90 Days", AMBER)
pagenum(s, 11)

# Left column — flagship feature
card(s, 0.28, 1.38, 7.6, 5.82, WHITE, accent_color=AMBER, accent_w=0.08)
rect(s, 0.44, 1.5, 2.4, 0.3, AMBER)
txt(s, "PRIORITY SHIP", 0.46, 1.52, 2.38, 0.26, size=9, bold=True, color=WHITE)
txt(s, "AI Semantic Search", 0.44, 1.9, 7.2, 0.5, size=20, bold=True, color=NEAR_BLACK)
txt(s, "The fastest credible AI feature TaskFlow can ship",
    0.44, 2.45, 7.2, 0.35, size=12, italic=True, color=GRAY_400)

rect(s, 0.44, 2.88, 7.25, 0.02, GRAY_200)

SEARCH_PTS = [
    ("What it does", "Users type natural language: 'overdue Q1 launch tasks' — AI returns ranked relevant results, not keyword matches"),
    ("Why this first", "Low hallucination risk (retrieval, not generation). Technically simpler than full chat. Ships in 4–6 weeks."),
    ("Impact", "Immediately addresses the 'chat with your tasks' use case. Credible demo for sales within weeks of launch."),
    ("Architecture", "Embed task data into a vector store via Anthropic/OpenAI API. No internal ML team required."),
]
for j, (heading, body) in enumerate(SEARCH_PTS):
    y = 3.05 + j * 0.98
    oval(s, 0.44, y + 0.1, 0.2, 0.2, AMBER)
    txt(s, heading, 0.74, y + 0.04, 2.0, 0.28, size=11, bold=True, color=AMBER)
    txt(s, body, 0.74, y + 0.34, 6.95, 0.58, size=11, color=GRAY_600)

# Right column — two cards stacked
for ci, (title, pts, clr) in enumerate([
    ("Publish Public AI Roadmap", [
        "Q2 2025 — AI Semantic Search  (In Development)",
        "Q3 2025 — AI Summaries & Status Updates",
        "Q4 2025 — Full Conversational AI Assistant",
        "Q1 2026 — Natural Language Automations",
    ], BLUE),
    ("Evaluate AI API Partnership", [
        "Anthropic / OpenAI API for Q&A layer — 8-12 weeks to v1",
        "Microsoft 365 Copilot connector for enterprise accounts",
        "Multi-model choice (GPT + Claude) as a differentiator",
        "API-first for v1; evaluate internal build for v2",
    ], NAVY),
]):
    y = 1.38 + ci * 2.96
    card(s, 8.08, y, 5.0, 2.78, WHITE, accent_color=clr, accent_w=0.08)
    txt(s, title, 8.26, y + 0.18, 4.65, 0.45, size=14, bold=True, color=NEAR_BLACK)
    rect(s, 8.26, y + 0.68, 4.65, 0.02, GRAY_200)
    for k, pt in enumerate(pts):
        oval(s, 8.28, y + 0.88 + k * 0.46, 0.13, 0.13, clr)
        txt(s, pt, 8.5, y + 0.82 + k * 0.46, 4.4, 0.44, size=11, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — LONG-TERM ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Full AI Vision for TaskFlow", "Long-term  ·  90+ Days", GREEN)
pagenum(s, 12)

INITIATIVES = [
    ("Full Conversational\nAI Assistant", GREEN,
     "Q4 2025",
     [
         "Ask: 'What's blocking the launch?', 'Which tasks am I behind on this week?'",
         "RAG architecture: index all tasks, projects, comments into a vector store",
         "Natural language task creation and updates via chat",
         "Slack / Teams integration for cross-platform context",
         "API-first v1 in 8-12 weeks; internal build for v2",
     ]),
    ("Natural Language\nAI Automations", BLUE,
     "Q1 2026",
     [
         "Describe automations in plain English — no builder required",
         "'When a task is overdue 3 days, alert assignee and escalate'",
         "Equivalent to Asana AI Studio and Monday AI Blocks",
         "This is the next major battleground after conversational AI",
     ]),
    ("Meeting Intelligence\nIntegration", NAVY,
     "Q1–Q2 2026",
     [
         "Connect Zoom / Google Meet transcripts to TaskFlow",
         "AI identifies action items, assigns owners, sets due dates",
         "Closes the loop: verbal commitment → TaskFlow task automatically",
         "Equivalent to ClickUp Notetaker & Notion Meeting Notes",
     ]),
]

for i, (title, clr, target, pts) in enumerate(INITIATIVES):
    x = 0.28 + i * 4.35
    # Header
    rect(s, x, 1.38, 4.15, 1.55, clr)
    txt(s, title, x + 0.2, 1.46, 3.75, 0.95,
        size=17, bold=True, color=WHITE)
    txt(s, f"Target: {target}", x + 0.2, 2.62, 3.75, 0.28,
        size=10, color=WHITE, italic=True)
    # Body card
    card(s, x, 2.98, 4.15, 4.22, WHITE)
    for j, pt in enumerate(pts):
        oy = 3.18 + j * 0.78
        if oy < 7.0:
            oval(s, x + 0.2, oy + 0.12, 0.18, 0.18, clr)
            txt(s, pt, x + 0.5, oy + 0.06, 3.5, 0.68, size=11.5, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PRICING OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "The AI Pricing Gap: TaskFlow's Competitive Wedge", "Strategic Opportunity", BLUE)
pagenum(s, 13)

txt(s, "Monthly AI cost for a 10-person team",
    0.28, 1.42, 8.5, 0.35, size=11, bold=True, color=GRAY_600)

COMPS = [
    ("ClickUp",  280, RED),
    ("Asana",    250, RED),
    ("Notion",   200, RED),
    ("Monday",   200, AMBER),
    ("Linear",   140, AMBER),
]
MAX_V = 280
BAR_MAX = 7.4

for i, (name, cost, clr) in enumerate(COMPS):
    y = 1.88 + i * 0.95
    bw = BAR_MAX * cost / MAX_V
    txt(s, name, 0.28, y + 0.12, 1.55, 0.58, size=13, bold=True, color=NEAR_BLACK)
    rect(s, 1.92, y + 0.14, bw, 0.58, clr)
    rect(s, 1.92, y + 0.14, bw, 0.58, clr)
    txt(s, f"~${cost}/mo", 1.98 + bw, y + 0.14, 1.3, 0.58,
        size=12, bold=True, color=clr)

# TaskFlow opportunity panel
rect(s, 9.28, 1.38, 3.8, 5.82, BLUE)
txt(s, "TaskFlow\nOpportunity", 9.42, 1.52, 3.5, 0.95,
    size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 9.42, 2.55, 3.5, 0.03, RGBColor(0x5B, 0x9B, 0xF0))

OPP_PTS = [
    "No competitor offers meaningful AI for free",
    "All charge $15–28/user/month on top of base plan",
    "Recommendation: bundle AI into existing paid plans",
    "Use AI as a retention and upgrade driver, not a surcharge",
    "Sales talking point: 'AI at no extra cost'",
]
for j, pt in enumerate(OPP_PTS):
    oval(s, 9.42, 2.72 + j * 0.84, 0.18, 0.18, RGBColor(0x5B, 0x9B, 0xF0))
    txt(s, pt, 9.7, 2.66 + j * 0.84, 3.2, 0.78, size=11.5, color=WHITE)

rect(s, 0.28, 6.65, 8.78, 0.52, BLUE_PALE, border=GRAY_200, border_w=0.4)
txt(s, "Key insight: every competitor has made AI an expensive surcharge. TaskFlow can own the "
       "narrative of 'AI that doesn't nickel-and-dime your team.'",
    0.42, 6.72, 8.55, 0.42, size=11, italic=True, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CEO TALKING POINTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "Talking Points for Sarah Chen, CEO", "Exec Messaging", NAVY)
pagenum(s, 14)

AUDIENCES = [
    ("Board & Investors", BLUE, BLUE_PALE,
     "The market has confirmed what we've believed: AI-powered task management is the future. "
     "Our competitors' launches validate the direction. We're positioned to respond with a "
     "focused, pragmatic roadmap — starting with shippable AI in Q2, scaling to full "
     "conversational AI by end of year."),
    ("Customer Conversations", GREEN, GREEN_PALE,
     "You've probably seen announcements from ClickUp, Notion, and Asana about AI features. "
     "We've been building our AI capabilities with the same care we brought to the rest of "
     "TaskFlow — focused on what actually helps you get work done, not AI for its own sake. "
     "Here's what's coming and when..."),
    ("Press & Analysts", AMBER, AMBER_PALE,
     "TaskFlow is building AI-native task management grounded in how teams actually work — "
     "not AI bolted onto legacy architecture. Our approach prioritizes retrieval accuracy, "
     "data privacy, and seamless workflow integration. Launching first AI features Q2 2025; "
     "full conversational assistant Q4."),
    ("Sales Team", NAVY, GRAY_100,
     "Yes, we're actively shipping AI features. Here's our roadmap: [show page]. "
     "Our competitors charge $200–280/month extra for a 10-person team just for AI. "
     "We're delivering AI value without punishing our customers with surcharges — "
     "that's the conversation to have with every prospect."),
]

for i, (audience, clr, pale, quote) in enumerate(AUDIENCES):
    x = 0.28 + (i % 2) * 6.55
    y = 1.38 + (i // 2) * 2.92
    card(s, x, y, 6.18, 2.78, pale)
    rect(s, x, y, 6.18, 0.5, clr)
    txt(s, f"For: {audience}", x + 0.18, y + 0.1, 5.8, 0.35,
        size=13, bold=True, color=WHITE)
    txt(s, f'"{quote}"', x + 0.18, y + 0.65, 5.82, 2.0,
        size=11.5, italic=True, color=NEAR_BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DECISIONS NEEDED
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "4 Decisions the Exec Team Must Make Now", "Leadership Decisions Required", BLUE)
pagenum(s, 15)

DECISIONS = [
    ("AI Pricing Model",
     "Include AI in paid plans  vs.  Charge as a separate add-on",
     "Include in paid plans. Use AI as a retention and upgrade driver. Directly counters competitors' 'AI tax' narrative."),
    ("Build vs. Integrate",
     "Internal AI engineering team  vs.  API partnership (Anthropic / OpenAI)",
     "API-first for v1 — delivers a compelling demo in 8–12 weeks vs. 6–12 months for internal build. Evaluate internal build for v2."),
    ("First Feature",
     "AI Semantic Search  vs.  AI Summaries  vs.  Full Conversational Chat",
     "Smart Search first — lowest hallucination risk, fastest to ship, directly addresses the 'chat with tasks' use case at manageable scope."),
    ("Public Roadmap",
     "Publish a public AI roadmap page  vs.  Keep roadmap internal",
     "Go public. Silence is worse than uncertainty. A visible roadmap gives sales reps something concrete to show prospects today."),
]

for i, (title, options, rec) in enumerate(DECISIONS):
    y = 1.38 + i * 1.44
    card(s, 0.28, y, 12.8, 1.36, WHITE)
    # Number badge
    rect(s, 0.28, y, 0.62, 1.36, BLUE)
    txt(s, str(i + 1), 0.28, y + 0.44, 0.62, 0.5,
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content
    txt(s, title, 1.05, y + 0.1, 11.7, 0.38, size=15, bold=True, color=NEAR_BLACK)
    txt(s, options, 1.05, y + 0.48, 11.7, 0.3,
        size=11, italic=True, color=GRAY_400)
    rect(s, 1.05, y + 0.82, 11.7, 0.02, GRAY_200)
    oval(s, 1.08, y + 0.94, 0.16, 0.16, GREEN)
    txt(s, rec, 1.32, y + 0.9, 11.42, 0.38, size=12, color=GRAY_600, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SUCCESS METRICS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
content_slide(s, "How We'll Know the Strategy Is Working", "Success Metrics", BLUE)
pagenum(s, 16)

METRICS = [
    ("AI-related churn rate",          "< 2% per quarter",      "90 days",  RED),
    ("Sales win rate (AI mentioned)",  "+10 pp vs. baseline",   "60 days",  AMBER),
    ("Smart Search: live in product",  "Shipped",               "60 days",  GREEN),
    ("Customer NPS post-AI launch",    "+8 points",             "90 days",  GREEN),
    ("Positive AI support tickets",    "50+ per month",         "90 days",  BLUE),
]

# Header
rect(s, 0.28, 1.38, 12.8, 0.52, NEAR_BLACK)
for label, cx, cw in [
    ("Metric",   0.5,  6.2),
    ("Target",   6.78, 3.6),
    ("Timeline", 10.5, 2.4),
]:
    txt(s, label, cx, 1.42, cw, 0.44, size=11, bold=True, color=WHITE)

for i, (metric, target, timeline, clr) in enumerate(METRICS):
    y = 1.9 + i * 1.0
    bg = WHITE if i % 2 == 0 else GRAY_100
    rect(s, 0.28, y, 12.8, 1.0, bg)
    # Color bar on left
    rect(s, 0.28, y, 0.08, 1.0, clr)
    txt(s, metric,   0.5,  y + 0.28, 6.1, 0.45, size=14, color=NEAR_BLACK)
    # Target pill
    rect(s, 6.78, y + 0.25, 3.55, 0.42, clr)
    txt(s, target, 6.84, y + 0.27, 3.43, 0.38,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, timeline, 10.5, y + 0.28, 2.4, 0.42, size=13, color=GRAY_600)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_slide(s, NEAR_BLACK)

# Left accent strip
rect(s, 0, 0, 0.08, 7.5, BLUE)

# Large decorative ">" mark
txt(s, ">", 9.5, 1.5, 3.5, 4.5, size=200, bold=True,
    color=RGBColor(0x1E, 0x30, 0x48), align=PP_ALIGN.CENTER)

txt(s, "The window is not closed.", 0.4, 1.2, 9.0, 1.0,
    size=44, bold=True, color=WHITE)
txt(s,
    "Every major competitor now has AI. But none has owned\n"
    "the narrative of AI that's powerful, affordable, and trusted.",
    0.4, 2.3, 9.0, 1.0, size=19, color=RGBColor(0x7B, 0xB3, 0xF0))

rect(s, 0.4, 3.42, 5.5, 0.04, BLUE)

txt(s, "Start this week:", 0.4, 3.62, 5, 0.38,
    size=14, bold=True, color=WHITE)

THIS_WEEK = [
    "Draft CEO email to customers  —  send by Friday",
    "Brief sales team with battle card  —  Monday standup",
    "Audit and remove 'no AI' messaging  —  this week",
    "Align exec team on pricing model and first feature  —  this week",
]
for j, item in enumerate(THIS_WEEK):
    oval(s, 0.42, 4.12 + j * 0.68, 0.2, 0.2, BLUE)
    txt(s, item, 0.72, 4.06 + j * 0.68, 8.3, 0.62, size=14, color=WHITE)

txt(s, "TaskFlow AI  ·  March 2025  ·  Confidential",
    0.4, 7.18, 9, 0.28, size=9, color=GRAY_400)
pagenum(s, 16, light=True)


# ── Save ───────────────────────────────────────────────────────────────────────
OUT = "/Users/stircrazy/Documents/claude-code-course/taskflow-ai-strategy.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
